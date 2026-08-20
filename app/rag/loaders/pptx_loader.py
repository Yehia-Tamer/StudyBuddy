from pptx import Presentation
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_core.documents import Document

class PPTXLoadError(Exception):
    pass

def extract_slide_text(slide) -> str:
    parts = []

    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text.strip()
            if text:
                parts.append(text)

    if slide.has_notes_slide and slide.notes_slide.notes_text_frame.text.strip():
        parts.append(f"[Speaker notes: {slide.notes_slide.notes_text_frame.text.strip()}]")

    return "\n".join(parts)

def load_pptx_slides(file_path: str) -> list[Document]:
    try:
        prs = Presentation(file_path)
    except Exception as e:
        raise PPTXLoadError(f"Could not open PowerPoint file: {e}")

    documents = []
    for i, slide in enumerate(prs.slides):
        text = extract_slide_text(slide)
        if text:
            documents.append(Document(
                page_content=text,
                metadata={"source_type": "pptx", "slide": i, "total_slides": len(prs.slides)}
            ))

    return documents

def load_and_split_pptx(file_path:str,chunk_size:int=800,chunk_overlap:int=150):
    documents=load_pptx_slides(file_path)
    if not documents:
        raise PPTXLoadError("No extractable text found in the PowerPoint file")

    splitter=RecursiveCharacterTextSplitter(
        chunk_size=chunk_size,
        chunk_overlap=chunk_overlap,
        separators=["\n\n", "\n", ". ", " ", ""]
    )

    return splitter.split_documents(documents)